#!/usr/bin/env python3
"""Generate the E2EE architecture overview deck (PowerPoint) with native shapes.

Run:  python docs/generate_deck.py
Out:  docs/E2EE-Architecture.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ─── Palette ────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
PANEL     = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT    = RGBColor(0x4C, 0x9A, 0xFF)   # blue
ACCENT2   = RGBColor(0x35, 0xD0, 0x9B)   # green
WARN      = RGBColor(0xF2, 0xB1, 0x4C)   # amber
DANGER    = RGBColor(0xE8, 0x6A, 0x6A)   # red
WHITE     = RGBColor(0xF5, 0xF8, 0xFF)
MUTED     = RGBColor(0xA9, 0xB8, 0xD6)
CLIENT    = RGBColor(0x2E, 0x5C, 0x9E)
SFUC      = RGBColor(0x6B, 0x47, 0x8F)
KDC       = RGBColor(0x2E, 0x7D, 0x6B)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def _set_text(tf, lines, size, color, bold=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, tuple):
            text, c, b, sz = ln
        else:
            text, c, b, sz = ln, color, bold, size
        run = p.add_run(); run.text = text
        run.font.size = Pt(sz); run.font.bold = b
        run.font.color.rgb = c; run.font.name = font


def text(s, x, y, w, h, lines, size=18, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    _set_text(tb.text_frame, lines, size, color, bold, align, anchor)
    return tb


def box(s, x, y, w, h, lines, fill=PANEL, line=ACCENT, size=14,
        color=WHITE, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        align=PP_ALIGN.CENTER, line_w=1.25):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    _set_text(sp.text_frame, lines, size, color, bold, align, MSO_ANCHOR.MIDDLE)
    sp.text_frame.margin_left = Inches(0.08)
    sp.text_frame.margin_right = Inches(0.08)
    sp.text_frame.margin_top = Inches(0.04)
    sp.text_frame.margin_bottom = Inches(0.04)
    return sp


def arrow(s, x1, y1, x2, y2, color=ACCENT, w=2.0, dashed=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    le = cn.line._get_or_add_ln()
    tail = le.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(tail)
    if dashed:
        d = le.makeelement(qn('a:prstDash'), {'val': 'dash'})
        le.insert(0, d)
    cn.shadow.inherit = False
    return cn


def header(s, title, kicker=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bar.shadow.inherit = False
    if kicker:
        text(s, Inches(0.6), Inches(0.30), Inches(11), Inches(0.4),
             kicker.upper(), size=13, color=ACCENT2, bold=True)
        ty = Inches(0.62)
    else:
        ty = Inches(0.45)
    text(s, Inches(0.6), ty, Inches(12.1), Inches(0.9), title,
         size=30, color=WHITE, bold=True)


def footer(s, n):
    text(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
         "E2EE WebRTC — Zero-Trust SFU (str0m + PERC)", size=10, color=MUTED)
    text(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
         str(n), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, items, size=18, gap=0.16):
    tb = s.shapes.add_textbox(x, y, w, Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap * 72)
        run = p.add_run()
        run.text = ("• " if lvl == 0 else "– ") + it
        run.font.size = Pt(size - lvl * 2)
        run.font.color.rgb = WHITE if lvl == 0 else MUTED
        run.font.name = "Segoe UI"
    return tb


N = 0
def num():
    global N; N += 1; return N


# ── Slide 1 — Title ─────────────────────────────────────────
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), EMU_W, Inches(2.6))
band.fill.solid(); band.fill.fore_color.rgb = PANEL; band.line.fill.background()
band.shadow.inherit = False
text(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.2),
     "End-to-End Encrypted WebRTC", size=46, color=WHITE, bold=True)
text(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
     "Zero-Trust SFU with PERC Double Encryption", size=26, color=ACCENT)
text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.6),
     "str0m (Rust SFU)  ·  Key Distributor (Node.js)  ·  Native libwebrtc Client",
     size=16, color=MUTED)
text(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.6),
     "Architecture Overview", size=14, color=ACCENT2, bold=True)

# ── Slide 2 — The Goal ──────────────────────────────────────
s = slide(); header(s, "The Goal: media the server can never read", "Problem")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.0), [
    "Multi-party WebRTC where the SFU routes media but cannot decrypt it.",
    "Compromising the server reveals only metadata — not audio or video.",
    "Inspired by IETF PERC: RFC 8871 (framework) + RFC 8723 (double encryption).",
    "Keys live only with clients and a trusted Key Distributor.",
])
box(s, Inches(7.1), Inches(1.9), Inches(5.4), Inches(1.1),
    ["Traditional SFU", "decrypts & re-encrypts media → server sees everything"],
    fill=PANEL, line=DANGER, size=15)
box(s, Inches(7.1), Inches(3.3), Inches(5.4), Inches(1.1),
    ["Zero-Trust SFU (this project)", "routes by RTP headers only → media stays sealed E2E"],
    fill=PANEL, line=ACCENT2, size=15)
box(s, Inches(7.1), Inches(4.7), Inches(5.4), Inches(1.3),
    ["Result", "Encrypted audio AND video verified end-to-end through the SFU"],
    fill=SFUC, line=ACCENT, size=16, bold=True)
footer(s, num())

# ── Slide 3 — Components ────────────────────────────────────
s = slide(); header(s, "Three components", "System")
box(s, Inches(0.7), Inches(1.9), Inches(3.8), Inches(2.4),
    ["SFU  (str0m, Rust)", "",
     "Hop-by-hop DTLS-SRTP per leg",
     "Routes by SSRC",
     "Forwards inner E2E payload",
     "Untrusted for secrecy"], fill=SFUC, line=ACCENT, size=14, bold=False)
box(s, Inches(4.8), Inches(1.9), Inches(3.8), Inches(2.4),
    ["Key Distributor (Node.js)", "",
     "Issues / rotates E2E keys",
     "Authenticates participants",
     "Never touches media",
     "Trusted"], fill=KDC, line=ACCENT2, size=14)
box(s, Inches(8.9), Inches(1.9), Inches(3.7), Inches(2.4),
    ["Native Client (C++/Node)", "",
     "libwebrtc + Node.js CLI",
     "Inner AES-128-GCM E2E",
     "Frame transformer",
     "Trusted"], fill=CLIENT, line=ACCENT, size=14)
box(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(1.5),
    ["Security property",
     "The SFU terminates only the hop-by-hop layer for routing. The inner end-to-end layer —",
     "and its key — are never available to the server. Compromise reveals metadata only."],
    fill=PANEL, line=ACCENT2, size=15, bold=False)
footer(s, num())

# ── Slide 4 — Status ────────────────────────────────────────
s = slide(); header(s, "Implementation status", "Where we are")
box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(1.2),
    ["Phase 1 — 1:1 DTLS-SRTP tunnel (SFU relays opaque packets)        ✅ Verified"],
    fill=PANEL, line=ACCENT, size=18, align=PP_ALIGN.LEFT)
box(s, Inches(0.7), Inches(3.3), Inches(11.9), Inches(1.2),
    ["Phase 2 — PERC double encryption (HBH SFU + frame E2E + Key Distributor)   ✅ Verified"],
    fill=PANEL, line=ACCENT2, size=18, align=PP_ALIGN.LEFT)
box(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.3),
    ["Now working end-to-end",
     "Two-way encrypted audio + video · keyframe relay · per-conference key rotation · unified config"],
    fill=SFUC, line=ACCENT, size=16, bold=True)
footer(s, num())

# ── Slide 5 — High-level architecture ───────────────────────
s = slide(); header(s, "High-level architecture", "Big picture")
ca = box(s, Inches(0.7), Inches(3.2), Inches(2.6), Inches(1.6),
         ["Client A", "encode → E2E → HBH"], fill=CLIENT, line=ACCENT, size=15, bold=True)
sfu = box(s, Inches(5.35), Inches(3.2), Inches(2.6), Inches(1.6),
          ["SFU (PERC)", "strip HBH · route · re-HBH"], fill=SFUC, line=ACCENT, size=15, bold=True)
cb = box(s, Inches(10.0), Inches(3.2), Inches(2.6), Inches(1.6),
         ["Client B", "HBH → E2E → decode"], fill=CLIENT, line=ACCENT, size=15, bold=True)
kd = box(s, Inches(5.35), Inches(1.5), Inches(2.6), Inches(1.0),
         ["Key Distributor", "E2E keys"], fill=KDC, line=ACCENT2, size=14, bold=True)
arrow(s, Inches(3.3), Inches(4.0), Inches(5.35), Inches(4.0), ACCENT, 2.5)
arrow(s, Inches(7.95), Inches(4.0), Inches(10.0), Inches(4.0), ACCENT, 2.5)
arrow(s, Inches(4.3), Inches(2.5), Inches(2.0), Inches(3.2), ACCENT2, 1.75, dashed=True)
arrow(s, Inches(9.0), Inches(2.5), Inches(11.3), Inches(3.2), ACCENT2, 1.75, dashed=True)
text(s, Inches(3.2), Inches(4.15), Inches(2.2), Inches(0.4),
     "HBH SRTP A", size=12, color=MUTED, align=PP_ALIGN.CENTER)
text(s, Inches(7.85), Inches(4.15), Inches(2.2), Inches(0.4),
     "HBH SRTP B", size=12, color=MUTED, align=PP_ALIGN.CENTER)
text(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.2),
     [("Inner E2E layer (green, dashed key delivery) is opaque to the SFU; "
       "only the hop-by-hop SRTP differs per leg.", MUTED, False, 15)])
footer(s, num())

# ── Slide 6 — Tunnel mode (Phase 1) ─────────────────────────
s = slide(); header(s, "Phase 1 — DTLS tunnel mode", "Foundation")
ca = box(s, Inches(0.7), Inches(2.6), Inches(2.6), Inches(1.5),
         ["Client A"], fill=CLIENT, line=ACCENT, size=16, bold=True)
sfu = box(s, Inches(5.35), Inches(2.6), Inches(2.6), Inches(1.5),
          ["SFU", "ICE relay only"], fill=SFUC, line=ACCENT, size=15, bold=True)
cb = box(s, Inches(10.0), Inches(2.6), Inches(2.6), Inches(1.5),
         ["Client B"], fill=CLIENT, line=ACCENT, size=16, bold=True)
arrow(s, Inches(3.3), Inches(3.35), Inches(5.35), Inches(3.35), MUTED, 2.2, dashed=True)
arrow(s, Inches(7.95), Inches(3.35), Inches(10.0), Inches(3.35), MUTED, 2.2, dashed=True)
text(s, Inches(2.9), Inches(2.0), Inches(7.5), Inches(0.5),
     "DTLS · SRTP · SRTCP forwarded as opaque bytes", size=14, color=MUTED,
     align=PP_ALIGN.CENTER)
bullets(s, Inches(0.7), Inches(4.5), Inches(11.8), [
    "SFU terminates only ICE/STUN; DTLS handshake runs end-to-end between clients.",
    "Fingerprint + SSRC swapping in the SDP answer makes E2E DTLS traverse the SFU.",
    "SFU has no SRTP keys at all — but it is point-to-point (1:1).",
])
footer(s, num())

# ── Slide 7 — PERC double encryption ────────────────────────
s = slide(); header(s, "Phase 2 — PERC double encryption", "Core design")
ca = box(s, Inches(0.7), Inches(3.0), Inches(2.7), Inches(2.0),
         ["Client A", "", "1  E2E encrypt", "2  HBH encrypt"],
         fill=CLIENT, line=ACCENT, size=14, bold=True)
sfu = box(s, Inches(5.3), Inches(3.0), Inches(2.7), Inches(2.0),
          ["SFU (PERC)", "", "strip HBH A", "read RTP headers", "re-HBH for B"],
          fill=SFUC, line=ACCENT, size=14, bold=True)
cb = box(s, Inches(9.9), Inches(3.0), Inches(2.7), Inches(2.0),
         ["Client B", "", "strip HBH B", "strip E2E", "decode"],
         fill=CLIENT, line=ACCENT, size=14, bold=True)
arrow(s, Inches(3.4), Inches(4.0), Inches(5.3), Inches(4.0), ACCENT, 2.5)
arrow(s, Inches(8.0), Inches(4.0), Inches(9.9), Inches(4.0), ACCENT, 2.5)
box(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.0),
    ["Outer = hop-by-hop SRTP (per leg, SFU has keys)   ·   "
     "Inner = end-to-end AES-128-GCM (SFU never has the key)"],
    fill=PANEL, line=ACCENT2, size=15, bold=True)
text(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.2),
     [("The SFU decrypts only the outer layer to read SSRC/PT/seq for routing, "
       "then re-encrypts the outer layer per receiver. The inner E2E payload is "
       "forwarded byte-for-byte.", MUTED, False, 15)])
footer(s, num())

# ── Slide 8 — Roles & trust ─────────────────────────────────
s = slide(); header(s, "Roles & trust boundaries", "Security model")
box(s, Inches(0.7), Inches(1.9), Inches(5.8), Inches(2.0),
    ["TRUSTED", "", "Client devices & app code",
     "libwebrtc (DTLS, SRTP, codecs)", "Key Distributor (holds E2E keys)"],
    fill=PANEL, line=ACCENT2, size=15, bold=False, align=PP_ALIGN.LEFT)
box(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(2.0),
    ["UNTRUSTED (for secrecy)", "", "SFU server & infrastructure",
     "Network between client & SFU", "Cloud provider / logs"],
    fill=PANEL, line=DANGER, size=15, bold=False, align=PP_ALIGN.LEFT)
box(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(1.8),
    ["SFU = partially trusted",
     "Trusted to relay packets (availability), NOT trusted with keys (secrecy).",
     "It can drop or delay packets, but can never read the E2E media."],
    fill=SFUC, line=ACCENT, size=16, bold=True)
footer(s, num())

# ── Slide 9 — Key distribution flow ─────────────────────────
s = slide(); header(s, "Key distribution flow", "Key Distributor")
steps = [
    "POST /conference — create",
    "POST /:id/join — KD mints E2E master key",
    "Key bundle returned (KEK, e2eMasterKey, kekSpi → key_id)",
    "WS /ws/endpoint — realtime key updates",
    "Member join/leave rotates KEK → 'rekey' pushed to all",
    "Client installs key: pc.installE2eeKey(keyId, keyBuf)",
]
y = 1.8
for i, st in enumerate(steps):
    box(s, Inches(0.7), Inches(y), Inches(0.7), Inches(0.6), [str(i + 1)],
        fill=ACCENT, line=ACCENT, size=16, bold=True)
    box(s, Inches(1.6), Inches(y), Inches(10.9), Inches(0.6), [st],
        fill=PANEL, line=ACCENT2 if i in (1, 4) else ACCENT, size=15,
        align=PP_ALIGN.LEFT)
    y += 0.78
footer(s, num())

# ── Slide 10 — Inner frame format ───────────────────────────
s = slide(); header(s, "Inner E2E frame format", "On the wire")
text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5),
     "AES-128-GCM at the encoded-frame boundary (libwebrtc FrameTransformer), empty AAD:",
     size=15, color=MUTED)
parts = [("key_id\n1B", ACCENT, 1.6), ("IV\n12B", ACCENT2, 2.4),
         ("ciphertext\nN bytes", SFUC, 5.0), ("GCM tag\n16B", WARN, 2.0)]
x = 0.9
for label, col, w in parts:
    box(s, Inches(x), Inches(2.5), Inches(w), Inches(1.1), [label],
        fill=PANEL, line=col, size=15, bold=True)
    x += w + 0.12
bullets(s, Inches(0.7), Inches(4.0), Inches(11.8), [
    "key_id — KEK SPI / epoch selector from the Key Distributor.",
    "IV — SSRC (4B, big-endian) ‖ frame counter (8B).",
    "Fixed overhead = 1 + 12 + 16 = 29 bytes per frame.",
    "SFU forwards this inner payload byte-for-byte (no per-packet rewriting).",
])
footer(s, num())

# ── Slide 11 — VP8 marker subtlety ──────────────────────────
s = slide(); header(s, "The VP8 keyframe marker", "Hard-won subtlety")
box(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(1.3),
    ["Problem",
     "Full-frame E2EE hides the codec bitstream. The depacketizer reads the VP8",
     "keyframe bit from byte 0 — now it's key_id. Every frame looks like a delta →",
     "decoder never starts → endless PLIs → black video."],
    fill=PANEL, line=DANGER, size=14, bold=False, align=PP_ALIGN.LEFT)
box(s, Inches(0.7), Inches(3.4), Inches(11.9), Inches(1.0),
    ["Fix — prepend a 1-byte cleartext marker to VIDEO frames",
     "0x00 = keyframe   ·   0x01 = delta   (from encoder IsKeyFrame())"],
    fill=PANEL, line=ACCENT2, size=15, bold=True)
# wire layout
box(s, Inches(0.9), Inches(4.7), Inches(1.4), Inches(0.9), ["marker\n1B"],
    fill=WARN, line=WARN, size=13, bold=True, color=BG)
box(s, Inches(2.4), Inches(4.7), Inches(9.6), Inches(0.9),
    ["key_id · IV · ciphertext · GCM tag  (inner E2E payload)"],
    fill=PANEL, line=ACCENT, size=14)
text(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.7),
     "Receiver strips the marker before the key-id check & decrypt; audio frames carry no marker.",
     size=14, color=MUTED)
footer(s, num())

# ── Slide 12 — Keyframe relay ───────────────────────────────
s = slide(); header(s, "Keyframe (PLI/FIR) relay", "RTCP across legs")
ca = box(s, Inches(0.7), Inches(3.0), Inches(2.7), Inches(1.4),
         ["Receiver B", "sends PLI"], fill=CLIENT, line=ACCENT, size=15, bold=True)
sfu = box(s, Inches(5.3), Inches(3.0), Inches(2.7), Inches(1.4),
          ["SFU", "map B → sender A"], fill=SFUC, line=ACCENT, size=15, bold=True)
cb = box(s, Inches(9.9), Inches(3.0), Inches(2.7), Inches(1.4),
         ["Sender A", "request_keyframe()"], fill=CLIENT, line=ACCENT, size=15, bold=True)
arrow(s, Inches(3.4), Inches(3.7), Inches(5.3), Inches(3.7), WARN, 2.5)
arrow(s, Inches(8.0), Inches(3.7), Inches(9.9), Inches(3.7), WARN, 2.5)
bullets(s, Inches(0.7), Inches(4.8), Inches(11.8), [
    "Because HBH SRTP terminates per leg, RTCP feedback terminates per leg too.",
    "On Event::KeyframeRequest the SFU maps the requester to the sending peer in the room.",
    "It calls request_keyframe(kind) on that sender's video rx stream — so a fresh keyframe flows.",
])
footer(s, num())

# ── Slide 13 — What the SFU sees ────────────────────────────
s = slide(); header(s, "What the SFU can — and cannot — see", "Security")
box(s, Inches(0.7), Inches(1.9), Inches(5.8), Inches(3.6),
    ["VISIBLE", "", "RTP headers (SSRC/PT/seq/ts)",
     "HBH SRTP for its own legs", "1-byte VP8 key/delta marker",
     "Packet sizes & timing"],
    fill=PANEL, line=WARN, size=16, bold=False, align=PP_ALIGN.LEFT)
box(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(3.6),
    ["HIDDEN", "", "E2E media payload (inner GCM)",
     "The E2E key (KD + clients only)", "Decoded media / codec bitstream",
     "Inner header extensions"],
    fill=PANEL, line=ACCENT2, size=16, bold=False, align=PP_ALIGN.LEFT)
text(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.7),
     "Server compromise leaks routing metadata only — never the media content.",
     size=15, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, num())

# ── Slide 14 — Configuration ────────────────────────────────
s = slide(); header(s, "Unified configuration system", "Operations")
box(s, Inches(0.7), Inches(1.8), Inches(3.6), Inches(1.1),
    ["config.json (JSONC)", "sectioned OR flat"], fill=PANEL, line=ACCENT,
    size=14, bold=True)
for i, (lbl, col) in enumerate([("sfu", ACCENT), ("keyDistributor", ACCENT2),
                                ("client", CLIENT)]):
    box(s, Inches(4.6 + i * 2.7), Inches(1.8), Inches(2.5), Inches(1.1),
        [lbl, "+ shared\nlogging/stats/diag"], fill=PANEL, line=col, size=13)
    arrow(s, Inches(4.3), Inches(2.35), Inches(4.6 + i * 2.7), Inches(2.35),
          MUTED, 1.5)
bullets(s, Inches(0.7), Inches(3.3), Inches(11.8), [
    "Resolution: --config (repeatable, deep-merge) → E2EE_CONFIG env → ./config.json.",
    "Same loader handles one combined file or per-host flat files; JSONC + trailing commas.",
    "Node apps share config-loader.js; str0m uses a matching loader in examples/util.",
    "Client pushes media params to native via env vars (width/height/fps/bitrate/codec).",
    "Toggles: log-to-file, periodic stats, wire log, per-frame E2E diagnostics.",
    "run-all.ps1 launches SFU + KD + two clients against the shared config.",
])
footer(s, num())

# ── Slide 15 — Roadmap ──────────────────────────────────────
s = slide(); header(s, "Done & what's next", "Roadmap")
box(s, Inches(0.7), Inches(1.9), Inches(5.8), Inches(3.4),
    ["COMPLETED ✅", "", "Key Distributor service",
     "Frame-level AES-128-GCM E2E", "PERC-capable native client",
     "VP8 keyframe marker", "Keyframe (PLI/FIR) relay", "Unified config + launcher"],
    fill=PANEL, line=ACCENT2, size=15, bold=False, align=PP_ALIGN.LEFT)
box(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(3.4),
    ["NEXT", "", "RFC 8723 at the SRTP layer",
     "Multi-party (N:N) routing", "EKT (RFC 8870) key transport",
     "Certificate pinning to identity", "Encrypted headers (Cryptex)",
     "MLS group key agreement"],
    fill=PANEL, line=ACCENT, size=15, bold=False, align=PP_ALIGN.LEFT)
footer(s, num())

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "E2EE-Architecture.pptx")
prs.save(out)
print("Wrote", out, "with", len(prs.slides._sldIdLst), "slides")
